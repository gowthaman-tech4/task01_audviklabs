"""
Document ingestion layer — handles text files, digital PDFs,
scanned/image PDFs (via OCR), and coordinate-aware extraction.
"""
import os

# Fix TESSDATA_PREFIX on Windows if it points to Tesseract-OCR instead of Tesseract-OCR\tessdata
tessdata_path = r"C:\Program Files\Tesseract-OCR\tessdata"
if os.path.exists(tessdata_path):
    os.environ["TESSDATA_PREFIX"] = tessdata_path

from dataclasses import dataclass, field
from typing import Optional, List, Dict, Tuple
from abc import ABC, abstractmethod


@dataclass
class WordBox:
    """A single word with its bounding box coordinates on a page image."""
    text: str
    x: int          # left pixel coordinate
    y: int          # top pixel coordinate
    w: int          # width in pixels
    h: int          # height in pixels
    page_num: int   # 0-indexed page number
    conf: int       # OCR confidence (0-100)


@dataclass
class CoordinateAwareResult:
    """Result from coordinate-aware ingestion."""
    full_text: str                          # Reconstructed plain text
    word_boxes: List[WordBox]               # Every word with coordinates
    char_to_boxes: Dict[int, List[WordBox]] # char index -> WordBox mapping
    page_images: list                       # PIL Image objects per page


def preprocess_image(pil_image) -> tuple:
    """
    Applies advanced OpenCV pre-processing to clean and enhance the image:
    1. Grayscale conversion.
    2. Rescaling/Upscaling (2x cubic interpolation) if resolution is low.
    3. Binarization (Otsu's thresholding) to eliminate noise.

    Returns:
        tuple of (processed_pil_image, scale_factor)
    """
    import cv2
    import numpy as np
    from PIL import Image

    img_np = np.array(pil_image)

    if len(img_np.shape) == 3:
        if img_np.shape[2] == 4:
            img_cv = cv2.cvtColor(img_np, cv2.COLOR_RGBA2BGR)
        elif img_np.shape[2] == 3:
            img_cv = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
        else:
            img_cv = img_np
    else:
        img_cv = img_np

    if len(img_cv.shape) == 3:
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
    else:
        gray = img_cv

    h, w = gray.shape[:2]
    scale_factor = 1.0
    if w < 2000 or h < 2000:
        gray = cv2.resize(gray, (w * 2, h * 2), interpolation=cv2.INTER_CUBIC)
        scale_factor = 2.0

    blurred = cv2.GaussianBlur(gray, (3, 3), 0)
    _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)

    return Image.fromarray(thresh), scale_factor


class BaseIngester(ABC):
    """Abstract base class for document ingesters."""

    @abstractmethod
    def extract_text(self, file_path: str) -> str:
        """Extract text from the given file."""
        pass

    @property
    @abstractmethod
    def supported_extensions(self) -> list:
        """Return list of supported file extensions."""
        pass


class TextIngester(BaseIngester):
    """Reads plain text files."""

    def extract_text(self, file_path: str) -> str:
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
        for enc in encodings:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        raise ValueError(f"Cannot decode file {file_path} with any supported encoding")

    @property
    def supported_extensions(self) -> list:
        return ['.txt', '.text', '.csv', '.log']


class PDFIngester(BaseIngester):
    """Extracts text from digital (selectable text) PDFs using pdfplumber."""

    def extract_text(self, file_path: str) -> str:
        try:
            import pdfplumber
        except ImportError:
            raise ImportError("pdfplumber is required: pip install pdfplumber")

        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        full_text = '\n\n'.join(text_parts)

        # If PDF has no selectable text, fall back to OCR
        if not full_text.strip():
            return OCRIngester().extract_text(file_path)

        return full_text

    @property
    def supported_extensions(self) -> list:
        return ['.pdf']


class OCRIngester(BaseIngester):
    """Extracts text from scanned/image PDFs and images using Tesseract OCR."""

    def __init__(self, tesseract_cmd: Optional[str] = None, lang: str = 'eng'):
        self.lang = lang
        if tesseract_cmd:
            try:
                import pytesseract
                pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
            except ImportError:
                pass

    def extract_text(self, file_path: str) -> str:
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError(
                "pytesseract and Pillow are required: "
                "pip install pytesseract Pillow"
            )

        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.pdf':
            return self._extract_from_pdf(file_path)
        elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'):
            return self._extract_from_image(file_path)
        else:
            raise ValueError(f"Unsupported file type for OCR: {ext}")

    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF by converting pages to images via PyMuPDF."""
        try:
            import fitz
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError(
                "PyMuPDF and pytesseract are required: pip install PyMuPDF pytesseract"
            )

        doc = fitz.open(file_path)
        text_parts = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=300)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            processed_img, _ = preprocess_image(img)
            page_text = pytesseract.image_to_string(processed_img, lang=self.lang)
            if page_text.strip():
                text_parts.append(page_text)
        doc.close()

        return '\n\n'.join(text_parts)

    def _extract_from_image(self, file_path: str) -> str:
        """Extract text from a single image."""
        import pytesseract
        from PIL import Image

        image = Image.open(file_path)
        processed_img, _ = preprocess_image(image)
        return pytesseract.image_to_string(processed_img, lang=self.lang)

    @property
    def supported_extensions(self) -> list:
        return ['.pdf', '.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif']


class CoordinateAwareIngester:
    """
    Extracts text WITH word-level bounding box coordinates.
    
    Converts every page to an image, runs Tesseract image_to_data(),
    and builds a character-to-coordinate mapping so detected PII
    entities can be projected back to exact pixel positions for redaction.
    """

    def __init__(self, dpi: int = 300, lang: str = 'eng'):
        self.dpi = dpi
        self.lang = lang

    def extract(self, file_path: str) -> CoordinateAwareResult:
        """
        Extract text and word-level coordinates from any supported file.
        
        Returns:
            CoordinateAwareResult with full_text, word_boxes, char mapping,
            and page images.
        """
        import pytesseract
        from PIL import Image

        ext = os.path.splitext(file_path)[1].lower()

        # Step 1: Convert to list of PIL images
        if ext == '.pdf':
            page_images = self._pdf_to_images(file_path)
        elif ext in ('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif'):
            page_images = [Image.open(file_path).convert("RGB")]
        else:
            # For text files, we can't do coordinate-aware extraction
            # Return a basic result
            text_ingester = TextIngester()
            text = text_ingester.extract_text(file_path)
            return CoordinateAwareResult(
                full_text=text,
                word_boxes=[],
                char_to_boxes={},
                page_images=[]
            )

        # Step 2: Run Tesseract image_to_data on each page
        all_word_boxes: List[WordBox] = []
        text_parts: List[str] = []

        for page_num, img in enumerate(page_images):
            processed_img, scale = preprocess_image(img)
            data = pytesseract.image_to_data(processed_img, lang=self.lang, output_type=pytesseract.Output.DICT)

            page_words = []
            n_boxes = len(data['text'])
            for i in range(n_boxes):
                word = data['text'][i].strip()
                conf = int(data['conf'][i])

                # Skip empty words and very low confidence
                if not word or conf < 10:
                    continue

                wb = WordBox(
                    text=word,
                    x=int(data['left'][i] / scale),
                    y=int(data['top'][i] / scale),
                    w=int(data['width'][i] / scale),
                    h=int(data['height'][i] / scale),
                    page_num=page_num,
                    conf=conf
                )
                all_word_boxes.append(wb)
                page_words.append(word)

            # Reconstruct page text from OCR words
            page_text = self._reconstruct_text(data)
            text_parts.append(page_text)

        # Step 3: Build the full text and character-to-box mapping
        full_text = '\n\n'.join(text_parts)
        char_to_boxes = self._build_char_mapping(full_text, all_word_boxes, text_parts)

        return CoordinateAwareResult(
            full_text=full_text,
            word_boxes=all_word_boxes,
            char_to_boxes=char_to_boxes,
            page_images=page_images
        )

    def _pdf_to_images(self, file_path: str) -> list:
        """Convert PDF pages to PIL images using PyMuPDF (fitz)."""
        import fitz
        from PIL import Image

        doc = fitz.open(file_path)
        images = []
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=self.dpi)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
        doc.close()
        return images

    def _reconstruct_text(self, data: dict) -> str:
        """
        Reconstruct readable text from Tesseract image_to_data output,
        preserving line breaks based on block/paragraph/line numbers.
        """
        lines = {}
        n_boxes = len(data['text'])

        for i in range(n_boxes):
            word = data['text'][i].strip()
            if not word:
                continue

            # Use (block_num, par_num, line_num) as line key
            line_key = (data['block_num'][i], data['par_num'][i], data['line_num'][i])
            if line_key not in lines:
                lines[line_key] = []
            lines[line_key].append(word)

        # Sort by line key and join
        sorted_keys = sorted(lines.keys())
        text_lines = []
        prev_block = None
        for key in sorted_keys:
            block_num = key[0]
            # Add extra newline between blocks
            if prev_block is not None and block_num != prev_block:
                text_lines.append('')
            prev_block = block_num
            text_lines.append(' '.join(lines[key]))

        return '\n'.join(text_lines)

    def _build_char_mapping(self, full_text: str, word_boxes: List[WordBox],
                            page_texts: List[str]) -> Dict[int, List[WordBox]]:
        """
        Build a mapping from character index in full_text to the WordBox(es)
        that contain that character.
        """
        char_to_boxes: Dict[int, List[WordBox]] = {}

        # We match word_boxes to their positions in full_text
        # by searching for each word sequentially
        search_start = 0
        for wb in word_boxes:
            # Find this word in the full text starting from search_start
            idx = full_text.find(wb.text, search_start)
            if idx == -1:
                # Try case-insensitive or partial match
                idx = full_text.lower().find(wb.text.lower(), search_start)

            if idx != -1:
                # Map every character in this word to this WordBox
                for ci in range(idx, idx + len(wb.text)):
                    if ci not in char_to_boxes:
                        char_to_boxes[ci] = []
                    char_to_boxes[ci].append(wb)
                # Advance search past this word
                search_start = idx + len(wb.text)

        return char_to_boxes


class DocxIngester(BaseIngester):
    """Extracts text from Microsoft Word (.docx) files."""

    def extract_text(self, file_path: str) -> str:
        try:
            import docx
        except ImportError:
            raise ImportError("python-docx is required: pip install python-docx")

        doc = docx.Document(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)

    @property
    def supported_extensions(self) -> list:
        return ['.docx']


class PptxIngester(BaseIngester):
    """Extracts text from Microsoft PowerPoint (.pptx) files."""

    def extract_text(self, file_path: str) -> str:
        try:
            import pptx
        except ImportError:
            raise ImportError("python-pptx is required: pip install python-pptx")

        prs = pptx.Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)

    @property
    def supported_extensions(self) -> list:
        return ['.pptx']


class XlsxIngester(BaseIngester):
    """Extracts text from Microsoft Excel (.xlsx) files."""

    def extract_text(self, file_path: str) -> str:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl is required: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet in wb.worksheets:
            text_parts.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell).strip() for cell in row if cell is not None]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)

    @property
    def supported_extensions(self) -> list:
        return ['.xlsx']


class DocumentIngester:
    """
    Unified document ingester that auto-detects file type and routes
    to the appropriate extractor.
    """

    def __init__(self):
        self._text_ingester = TextIngester()
        self._pdf_ingester = PDFIngester()
        self._ocr_ingester = OCRIngester()
        self._coord_ingester = CoordinateAwareIngester()
        self._docx_ingester = DocxIngester()
        self._pptx_ingester = PptxIngester()
        self._xlsx_ingester = XlsxIngester()

    def extract_text(self, file_path: str, force_ocr: bool = False) -> str:
        """
        Extract text from any supported document type.

        Args:
            file_path: Path to the document file.
            force_ocr: If True, always use OCR even for digital PDFs.

        Returns:
            Extracted text content.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower()

        if force_ocr:
            return self._ocr_ingester.extract_text(file_path)

        if ext in self._text_ingester.supported_extensions:
            return self._text_ingester.extract_text(file_path)
        elif ext in self._docx_ingester.supported_extensions:
            return self._docx_ingester.extract_text(file_path)
        elif ext in self._xlsx_ingester.supported_extensions:
            return self._xlsx_ingester.extract_text(file_path)
        elif ext in self._pptx_ingester.supported_extensions:
            return self._pptx_ingester.extract_text(file_path)
        elif ext == '.pdf':
            return self._pdf_ingester.extract_text(file_path)
        elif ext in self._ocr_ingester.supported_extensions:
            return self._ocr_ingester.extract_text(file_path)
        else:
            # Try as text file
            try:
                return self._text_ingester.extract_text(file_path)
            except Exception:
                raise ValueError(
                    f"Unsupported file type: {ext}. "
                    f"Supported: .txt, .pdf, .png, .jpg, .jpeg, .tiff, .bmp"
                )

    def extract_with_coordinates(self, file_path: str) -> CoordinateAwareResult:
        """
        Extract text WITH word-level bounding box coordinates.
        Use this for coordinate-based redaction.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        return self._coord_ingester.extract(file_path)
